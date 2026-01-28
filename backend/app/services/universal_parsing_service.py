from __future__ import annotations

import json
import re
import tempfile
import logging
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import httpx

import pandas as pd
from docx import Document
from pptx import Presentation

from app.core.config import settings
from app.models.score import ScoreItem, StudentScore
from app.services.azure_openai_responses_client import AzureOpenAIResponsesClient


logger = logging.getLogger(__name__)


@dataclass(frozen=True)
class ParsePreview:
    file_type: str
    ir: dict[str, Any]
    preview: dict[str, Any]


@dataclass(frozen=True)
class MappingResult:
    mapping: dict[str, Any]
    confidence: float
    errors: list[str]
    recommendations: list[str]
    usage: dict[str, int]
    raw_text: str


def _resolve_responses_url() -> str:
    if settings.AZURE_OPENAI_RESPONSES_URL and settings.AZURE_OPENAI_RESPONSES_URL.strip():
        return settings.AZURE_OPENAI_RESPONSES_URL.strip()

    if not settings.AZURE_OPENAI_ENDPOINT or not settings.AZURE_OPENAI_ENDPOINT.strip():
        raise ValueError("AZURE_OPENAI_ENDPOINT or AZURE_OPENAI_RESPONSES_URL must be set")

    endpoint = settings.AZURE_OPENAI_ENDPOINT.strip().rstrip("/")
    if endpoint.endswith("/openai/v1"):
        return f"{endpoint}/responses"
    if endpoint.endswith("/openai/v1/responses"):
        return endpoint
    return f"{endpoint}/openai/v1/responses"


def _resolve_responses_url_2() -> str | None:
    v = getattr(settings, "AZURE_OPENAI_RESPONSES_URL_2", None)
    if v and str(v).strip():
        return str(v).strip().rstrip("/")
    return None


def _extract_json_object(text: str) -> dict[str, Any]:
    text = (text or "").strip()
    if not text:
        raise ValueError("Empty model output")

    # Fast path
    try:
        obj = json.loads(text)
        if isinstance(obj, dict):
            return obj
    except Exception:
        pass

    # Heuristic: find first {...} block
    start = text.find("{")
    end = text.rfind("}")
    if start >= 0 and end > start:
        candidate = text[start : end + 1]
        obj = json.loads(candidate)
        if isinstance(obj, dict):
            return obj

    raise ValueError("Model output is not a JSON object")


def _safe_float(v: Any, default: float = 0.0) -> float:
    try:
        if v is None:
            return default
        if isinstance(v, (int, float)):
            return float(v)
        s = str(v).strip()
        if not s:
            return default
        return float(s)
    except Exception:
        return default


class UniversalParsingService:
    """Universal parsing: Extract IR -> infer mapping via AI -> deterministic parse."""

    @staticmethod
    def _resolve_parsing_model() -> str:
        if settings.PARSING_MODEL and settings.PARSING_MODEL.strip():
            return settings.PARSING_MODEL.strip()
        # Fallback to analysis model / deployment name for convenience.
        if settings.ANALYSIS_MODEL and settings.ANALYSIS_MODEL.strip():
            return settings.ANALYSIS_MODEL.strip()
        if settings.AZURE_OPENAI_DEPLOYMENT_NAME and settings.AZURE_OPENAI_DEPLOYMENT_NAME.strip():
            return settings.AZURE_OPENAI_DEPLOYMENT_NAME.strip()
        raise ValueError("PARSING_MODEL (or ANALYSIS_MODEL / AZURE_OPENAI_DEPLOYMENT_NAME) must be set")

    @staticmethod
    def extract_preview(*, file_bytes: bytes, filename: str) -> ParsePreview:
        suffix = Path(filename).suffix.lower()
        if suffix not in (".xlsx", ".docx", ".pptx"):
            raise ValueError("Unsupported file type")

        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tf:
            tf.write(file_bytes)
            temp_path = tf.name

        try:
            if suffix == ".xlsx":
                ir, preview = _extract_excel_ir_and_preview(temp_path)
                return ParsePreview(file_type="xlsx", ir=ir, preview=preview)
            if suffix == ".docx":
                ir, preview = _extract_word_ir_and_preview(temp_path)
                return ParsePreview(file_type="docx", ir=ir, preview=preview)
            ir, preview = _extract_ppt_ir_and_preview(temp_path)
            return ParsePreview(file_type="pptx", ir=ir, preview=preview)
        finally:
            try:
                Path(temp_path).unlink(missing_ok=True)
            except Exception:
                pass

    @staticmethod
    async def infer_mapping(*, file_type: str, ir: dict[str, Any], preview: dict[str, Any]) -> MappingResult:
        if not settings.AZURE_OPENAI_API_KEY:
            raise ValueError("AZURE_OPENAI_API_KEY must be set")
        parsing_model = UniversalParsingService._resolve_parsing_model()

        client = AzureOpenAIResponsesClient(
            responses_url=_resolve_responses_url(),
            api_key=settings.AZURE_OPENAI_API_KEY,
            fallback_responses_url=_resolve_responses_url_2(),
            fallback_api_key=(settings.AZURE_OPENAI_API_KEY_2 or None),
            timeout_seconds=float(settings.OPENAI_REQUEST_TIMEOUT_SECONDS or 600.0),
        )

        system_prompt = (
            "You are a strict JSON generator. Output ONLY valid JSON (no markdown, no comments). "
            "Your task: infer a parsing mapping plan to convert the given file IR into a list of StudentScore objects. "
            "If uncertain, set confidence low and include errors/recommendations."
        )

        # Responses structured output: prefer json_schema for strictness.
        mapping_schema: dict[str, Any] = {
            "type": "object",
            "additionalProperties": False,
            "properties": {
                "confidence": {"type": "number", "minimum": 0, "maximum": 1},
                "mapping": {"type": "object", "additionalProperties": True},
                "errors": {"type": "array", "items": {"type": "string"}},
                "recommendations": {"type": "array", "items": {"type": "string"}},
            },
            "required": ["confidence", "mapping", "errors", "recommendations"],
        }

        user_prompt = json.dumps(
            {
                "task": "infer_mapping_plan",
                "file_type": file_type,
                "ir": ir,
                "preview": preview,
                "output_schema": {
                    "confidence": "number 0..1",
                    "mapping": "object (see below)",
                    "errors": "string[]",
                    "recommendations": "string[]",
                    "mapping_schema": {
                        "common": {
                            "student_name": {
                                "source": "excel_column|text_heading|slide_title",
                                "column": "string or int (excel)",
                                "row_start": "int (excel, data start row index, default 1)",
                            },
                            "total_score": {
                                "column": "string or int (excel, optional)",
                            },
                            "items": {
                                "mode": "marker|explicit",
                                "default_deduction": "number (for marker mode)",
                                "columns": "(excel) list of column names or indices",
                                "line_pattern": "(docx/pptx) regex like '^(.+?)[:：]\s*(\\d+(?:\\.\\d+)?)$'",
                            },
                        },
                        "excel": {
                            "sheet": "string (preferred) or null",
                            "header_row": "int (0-based)",
                            "data_start_row": "int (0-based, usually header_row+1)",
                        },
                    },
                },
            },
            ensure_ascii=False,
        )

        try:
            result = await client.create_text_response(
                model=parsing_model,
                fallback_model=(settings.PARSING_MODEL_2.strip() if settings.PARSING_MODEL_2 else None),
                system_prompt=system_prompt,
                user_prompt=user_prompt,
                reasoning_effort=(settings.PARSING_REASONING_EFFORT or "high"),
                text_format={
                    "type": "json_schema",
                    "name": "mapping_plan",
                    "schema": mapping_schema,
                    "strict": True,
                },
            )
        except httpx.HTTPStatusError as e:
            # Some deployments/models may not support json_schema yet; fallback to json_object.
            status_code = getattr(getattr(e, "response", None), "status_code", None)
            if status_code in (400, 404, 422):
                logger.warning("AOAI structured output json_schema rejected (HTTP %s). Fallback to json_object.", status_code)
                result = await client.create_text_response(
                    model=parsing_model,
                    fallback_model=(settings.PARSING_MODEL_2.strip() if settings.PARSING_MODEL_2 else None),
                    system_prompt=system_prompt,
                    user_prompt=user_prompt,
                    reasoning_effort=(settings.PARSING_REASONING_EFFORT or "high"),
                    text_format={"type": "json_object"},
                )
            else:
                raise

        raw_text = result.text
        obj = _extract_json_object(raw_text)

        confidence = _safe_float(obj.get("confidence"), default=0.0)
        mapping = obj.get("mapping") if isinstance(obj.get("mapping"), dict) else {}
        errors = obj.get("errors") if isinstance(obj.get("errors"), list) else []
        recommendations = obj.get("recommendations") if isinstance(obj.get("recommendations"), list) else []

        usage = {"prompt_tokens": int(result.usage.input_tokens), "completion_tokens": int(result.usage.output_tokens)}

        return MappingResult(
            mapping=mapping,
            confidence=confidence,
            errors=[str(e) for e in errors if str(e).strip()],
            recommendations=[str(r) for r in recommendations if str(r).strip()],
            usage=usage,
            raw_text=raw_text,
        )

    @staticmethod
    def parse_full(*, file_bytes: bytes, filename: str, mapping: dict[str, Any]) -> List[StudentScore]:
        suffix = Path(filename).suffix.lower()
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tf:
            tf.write(file_bytes)
            temp_path = tf.name

        try:
            if suffix == ".xlsx":
                return _parse_excel_full(temp_path, mapping)
            if suffix == ".docx":
                return _parse_word_full(temp_path, mapping)
            if suffix == ".pptx":
                return _parse_ppt_full(temp_path, mapping)
            raise ValueError("Unsupported file type")
        finally:
            try:
                Path(temp_path).unlink(missing_ok=True)
            except Exception:
                pass


def _extract_excel_ir_and_preview(file_path: str) -> Tuple[dict[str, Any], dict[str, Any]]:
    # Read a limited amount for IR/preview to keep memory bounded.
    sheets = pd.read_excel(file_path, sheet_name=None, header=None, nrows=200)
    sheet_names = list(sheets.keys())

    # Default to first sheet for preview
    first_name = sheet_names[0] if sheet_names else None
    first_df = sheets[first_name] if first_name else pd.DataFrame()

    # Limit for IR/preview
    sample_df = first_df.iloc[:40, :40].copy() if not first_df.empty else first_df

    sample_rows: list[list[Any]] = []
    normalized_df = sample_df.copy()
    for r_idx, row in sample_df.iterrows():
        norm_row = [_normalize_cell(v) for v in row.tolist()]
        sample_rows.append(norm_row)
        normalized_df.loc[r_idx, :] = norm_row

    # Detect header row candidates from the first few rows.
    header_candidates: list[dict[str, Any]] = []
    max_check_rows = min(8, len(sample_rows))
    col_count = int(sample_df.shape[1]) if not sample_df.empty else 0

    def is_text_like(x: Any) -> bool:
        if x is None:
            return False
        if isinstance(x, (int, float)):
            return False
        s = str(x).strip()
        if not s:
            return False
        # avoid treating pure numbers as headers
        return re.fullmatch(r"[-+]?\d+(?:\.\d+)?", s) is None

    for i in range(max_check_rows):
        row = sample_rows[i]
        non_null = [c for c in row if c is not None]
        if len(non_null) < 3:
            continue
        text_cells = [c for c in row if is_text_like(c)]
        if not text_cells:
            continue
        unique_text = len({str(c).strip() for c in text_cells})
        text_ratio = len(text_cells) / max(1, len(non_null))
        unique_ratio = unique_text / max(1, len(text_cells))
        density = len(non_null) / max(1, col_count)
        score = round((text_ratio * 0.6 + unique_ratio * 0.3 + density * 0.1), 4)
        header_candidates.append(
            {
                "row_index": i,
                "score": score,
                "non_null": len(non_null),
                "text_ratio": round(text_ratio, 4),
                "unique_text_ratio": round(unique_ratio, 4),
            }
        )

    header_candidates.sort(key=lambda x: x.get("score", 0), reverse=True)
    header_candidates = header_candidates[:3]

    suggested_header_row = header_candidates[0]["row_index"] if header_candidates else 0

    # Column labels derived from suggested header row if possible.
    column_labels: list[str] = []
    if sample_rows and 0 <= suggested_header_row < len(sample_rows):
        hr = sample_rows[suggested_header_row]
        for idx, cell in enumerate(hr):
            label = str(cell).strip() if is_text_like(cell) else ""
            column_labels.append(label or f"#{idx + 1}")
    else:
        column_labels = [f"#{i + 1}" for i in range(col_count)]

    # Per-column stats (from the limited sample)
    column_stats: list[dict[str, Any]] = []
    if col_count:
        for c in range(col_count):
            col_vals = [row[c] if c < len(row) else None for row in sample_rows[: max(1, max_check_rows)]]
            non_null_vals = [v for v in col_vals if v is not None]
            numeric_vals: list[float] = []
            text_lens: list[int] = []
            for v in non_null_vals:
                if isinstance(v, (int, float)):
                    numeric_vals.append(float(v))
                else:
                    text_lens.append(len(str(v)))
            column_stats.append(
                {
                    "index": c,
                    "label": column_labels[c] if c < len(column_labels) else f"#{c + 1}",
                    "null_ratio": round(1 - (len(non_null_vals) / max(1, len(col_vals))), 4),
                    "numeric_ratio": round(len(numeric_vals) / max(1, len(non_null_vals)), 4) if non_null_vals else 0.0,
                    "avg_text_len": round(sum(text_lens) / max(1, len(text_lens)), 2) if text_lens else 0.0,
                }
            )

    ir = {
        "sheet_names": sheet_names,
        "first_sheet": first_name,
        "shape": {"rows": int(first_df.shape[0]), "cols": int(first_df.shape[1])},
        "header_row_candidates": header_candidates,
        "suggested_header_row": suggested_header_row,
        "column_stats": column_stats,
    }

    preview = {
        "sheet": first_name,
        "sample_rows": sample_rows[:15],
        "sample_row_count": len(sample_rows[:15]),
        "sample_col_count": len(sample_rows[0]) if sample_rows else 0,
        "column_labels": column_labels[: (len(sample_rows[0]) if sample_rows else 0)],
        "suggested_header_row": suggested_header_row,
    }

    return ir, preview


def _extract_word_ir_and_preview(file_path: str) -> Tuple[dict[str, Any], dict[str, Any]]:
    doc = Document(file_path)

    paras = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    paras_sample = paras[:40]

    tables_sample: list[list[list[str]]] = []
    for tbl in doc.tables[:3]:
        grid: list[list[str]] = []
        for r in tbl.rows[:10]:
            grid.append([c.text.strip() for c in r.cells[:10]])
        tables_sample.append(grid)

    ir = {
        "paragraph_count": len(paras),
        "table_count": len(doc.tables),
    }

    preview = {
        "paragraphs": paras_sample,
        "tables": tables_sample,
    }

    return ir, preview


def _extract_ppt_ir_and_preview(file_path: str) -> Tuple[dict[str, Any], dict[str, Any]]:
    prs = Presentation(file_path)

    slide_previews: list[dict[str, Any]] = []
    for i, slide in enumerate(prs.slides[:10]):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    texts.append(t)
        slide_previews.append({"slide_index": i, "texts": texts[:20]})

    ir = {
        "slide_count": len(prs.slides),
    }

    preview = {
        "slides": slide_previews,
    }

    return ir, preview


def _normalize_cell(v: Any) -> Any:
    if v is None:
        return None
    try:
        if pd.isna(v):
            return None
    except Exception:
        pass
    if isinstance(v, float) and (v != v):
        return None
    if isinstance(v, (int, float)):
        return v
    s = str(v).strip()
    return s if s else None


def _resolve_excel_column(df: pd.DataFrame, col: Any) -> str:
    if isinstance(col, int):
        return df.columns[int(col)]
    if isinstance(col, str):
        # Exact match first
        if col in df.columns:
            return col
        # Try stringified index
        if col.isdigit():
            idx = int(col)
            return df.columns[idx]
        # Fuzzy contains
        for c in df.columns:
            if str(c).strip() == col.strip():
                return c
    raise ValueError(f"Unknown column: {col}")


def _auto_detect_total_score_column(df: pd.DataFrame) -> Optional[str]:
    """Best-effort detection of a total score column.

    Some files have a clear '总分' column, but the AI mapping may omit it.
    If we fail to detect, callers will fall back to default score.
    """

    if df is None or df.empty or len(df.columns) == 0:
        return None

    keywords = ["总分", "得分", "总成绩", "总计", "总评", "score", "total"]
    for c in df.columns:
        if not isinstance(c, str):
            continue
        cc = c.strip().lower()
        for kw in keywords:
            if kw in cc:
                return c

    # Heuristic: last column is mostly numeric.
    last = df.columns[-1]
    try:
        s = pd.to_numeric(df[last], errors="coerce")
        if len(s) > 0 and float(s.notna().mean()) >= 0.6:
            return str(last)
    except Exception:
        return None
    return None


def _parse_excel_full(file_path: str, mapping: dict[str, Any]) -> List[StudentScore]:
    excel_cfg = mapping.get("excel") if isinstance(mapping.get("excel"), dict) else {}
    sheet = excel_cfg.get("sheet")
    header_row = int(excel_cfg.get("header_row", 0))

    df = pd.read_excel(file_path, sheet_name=sheet or 0, header=header_row)

    common = mapping.get("common") if isinstance(mapping.get("common"), dict) else mapping

    student_name_cfg = common.get("student_name") if isinstance(common.get("student_name"), dict) else {}
    student_col = student_name_cfg.get("column", 0)

    total_cfg = common.get("total_score") if isinstance(common.get("total_score"), dict) else {}
    total_col = total_cfg.get("column")

    items_cfg = common.get("items") if isinstance(common.get("items"), dict) else {}
    mode = str(items_cfg.get("mode") or "marker").lower()
    default_deduction = _safe_float(items_cfg.get("default_deduction"), default=1.0)

    item_cols = items_cfg.get("columns")
    if not isinstance(item_cols, list) or not item_cols:
        # Fallback: all columns except name and total
        exclude = set()
        try:
            exclude.add(_resolve_excel_column(df, student_col))
        except Exception:
            pass
        if total_col is not None:
            try:
                exclude.add(_resolve_excel_column(df, total_col))
            except Exception:
                pass
        item_cols = [c for c in df.columns if c not in exclude]

    try:
        student_col_name = _resolve_excel_column(df, student_col)
    except Exception:
        student_col_name = df.columns[0]

    total_col_name: Optional[str] = None
    if total_col is not None:
        try:
            total_col_name = _resolve_excel_column(df, total_col)
        except Exception:
            total_col_name = None
    else:
        total_col_name = _auto_detect_total_score_column(df)

    resolved_item_cols: list[str] = []
    for c in item_cols:
        try:
            resolved_item_cols.append(_resolve_excel_column(df, c))
        except Exception:
            continue

    # Defensive: never treat name/total columns as item columns.
    resolved_item_cols = [
        c for c in resolved_item_cols if c != student_col_name and (total_col_name is None or c != total_col_name)
    ]

    scores: list[StudentScore] = []

    for _, row in df.iterrows():
        student_name = str(row.get(student_col_name, "")).strip()
        if not student_name or student_name.lower() == "nan":
            continue

        total_score = 100.0
        if total_col_name:
            total_score = _safe_float(row.get(total_col_name), default=100.0)

        items: list[ScoreItem] = []
        for c in resolved_item_cols:
            val = row.get(c)
            if mode == "marker":
                if val is None:
                    continue
                try:
                    if pd.isna(val):
                        continue
                except Exception:
                    pass
                if str(val).strip() == "" or val == 0:
                    continue
                items.append(
                    ScoreItem(
                        question_name=str(c),
                        deduction=float(default_deduction),
                        category=_guess_category(str(c)),
                    )
                )
            else:
                # explicit numeric values
                dv = _safe_float(val, default=0.0)
                if dv == 0.0:
                    continue
                items.append(
                    ScoreItem(
                        question_name=str(c),
                        deduction=float(abs(dv)),
                        category=_guess_category(str(c)),
                    )
                )

        scores.append(StudentScore(student_name=student_name, scores=items, total_score=float(total_score)))

    return scores


def _parse_word_full(file_path: str, mapping: dict[str, Any]) -> List[StudentScore]:
    doc = Document(file_path)
    lines = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]

    common = mapping.get("common") if isinstance(mapping.get("common"), dict) else mapping
    items_cfg = common.get("items") if isinstance(common.get("items"), dict) else {}
    pattern = items_cfg.get("line_pattern") or r"^(.+?)[:：]\s*(\d+(?:\.\d+)?)$"
    rx = re.compile(pattern)

    students: list[StudentScore] = []
    current_name: Optional[str] = None
    current_items: list[ScoreItem] = []

    def flush():
        nonlocal current_name, current_items
        if current_name:
            total = sum(i.deduction for i in current_items) if current_items else 100.0
            students.append(StudentScore(student_name=current_name, scores=current_items, total_score=float(total)))
        current_name = None
        current_items = []

    for line in lines:
        # Student name heuristic: short line without ':'
        if (":" not in line and "：" not in line) and len(line.split()) <= 2 and len(line) <= 20:
            # Start new student section
            if current_name is not None:
                flush()
            current_name = line
            continue

        m = rx.match(line)
        if not m or not current_name:
            continue

        q = m.group(1).strip()
        d = _safe_float(m.group(2), default=0.0)
        if d == 0.0:
            continue
        current_items.append(ScoreItem(question_name=q, deduction=float(abs(d)), category=_guess_category(q)))

    if current_name is not None:
        flush()

    return students


def _parse_ppt_full(file_path: str, mapping: dict[str, Any]) -> List[StudentScore]:
    prs = Presentation(file_path)

    common = mapping.get("common") if isinstance(mapping.get("common"), dict) else mapping
    items_cfg = common.get("items") if isinstance(common.get("items"), dict) else {}
    pattern = items_cfg.get("line_pattern") or r"^(.+?)[:：]\s*(\d+(?:\.\d+)?)$"
    rx = re.compile(pattern)

    students: list[StudentScore] = []

    for slide in prs.slides:
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    texts.extend([x.strip() for x in t.split("\n") if x.strip()])

        if not texts:
            continue

        # Student name heuristic: first short text
        student_name = None
        for t in texts:
            if (":" not in t and "：" not in t) and len(t.split()) <= 2 and len(t) <= 20:
                student_name = t
                break
        if not student_name:
            continue

        items: list[ScoreItem] = []
        for t in texts:
            m = rx.match(t)
            if not m:
                continue
            q = m.group(1).strip()
            d = _safe_float(m.group(2), default=0.0)
            if d == 0.0:
                continue
            items.append(ScoreItem(question_name=q, deduction=float(abs(d)), category=_guess_category(q)))

        total = sum(i.deduction for i in items) if items else 100.0
        students.append(StudentScore(student_name=student_name, scores=items, total_score=float(total)))

    return students


def _guess_category(question_name: str) -> str:
    q = (question_name or "").lower()
    if "选择" in q or "单选" in q or "多选" in q:
        return "选择题"
    if "填空" in q:
        return "填空题"
    if "计算" in q:
        return "计算题"
    if "应用" in q or "解答" in q:
        return "应用题"
    if "判断" in q:
        return "判断题"
    if "简答" in q:
        return "简答题"
    if "作图" in q or "画图" in q:
        return "作图题"
    return question_name
