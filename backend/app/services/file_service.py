import pandas as pd
from docx import Document
from pptx import Presentation
from typing import List, Dict, Optional
import os
import logging
from app.models.score import StudentScore, ScoreItem

# 配置日志
logger = logging.getLogger(__name__)

class FileService:
    @staticmethod
    async def process_excel(file_path: str) -> List[StudentScore]:
        """
        处理Excel文件并提取成绩信息
        新格式：第一行为知识点，第一列（从第二行开始）为学生姓名，
        单元格有值表示该学生在该知识点有扣分（仅作标记，不代表真实扣分），
        空值表示不扣分，最后一列为总分
        """
        try:
            logger.info(f"开始读取Excel文件: {file_path}")
            # 读取Excel文件，第一行作为表头
            df = pd.read_excel(file_path, header=0)
            logger.info(f"Excel文件读取成功，共 {len(df)} 行, {len(df.columns)} 列")
            
            student_scores = []
            
            # 获取知识点列表（第一行，从第二列到倒数第二列，最后一列是总分）
            knowledge_points = df.columns[1:-1].tolist()
            logger.info(f"识别到 {len(knowledge_points)} 个知识点")
            
            # 遍历每一行（每一行代表一个学生）
            for index, row in df.iterrows():
                try:
                    # 获取学生姓名（第一列）
                    student_name = str(row.iloc[0]).strip()
                    
                    # 跳过空的或无效的学生姓名
                    if pd.isna(student_name) or student_name == '' or student_name.lower() == 'nan':
                        logger.warning(f"跳过无效的学生名称: {student_name}")
                        continue
                    
                    logger.info(f"处理学生: {student_name}")
                    
                    # 获取总分（最后一列）
                    try:
                        total_score = float(row.iloc[-1])
                        logger.info(f"{student_name} 总分: {total_score}")
                    except (ValueError, TypeError) as e:
                        logger.error(f"无法读取 {student_name} 的总分: {str(e)}")
                        total_score = 100.0  # 默认满分
                    
                    # 收集所有有标记的知识点（扣分项）
                    scores = []
                    
                    # 遍历每个知识点列（从第二列到倒数第二列）
                    for col_idx, knowledge_point in enumerate(knowledge_points, start=1):
                        value = row.iloc[col_idx]
                        
                        # 如果单元格有值（不为空），表示这是一个扣分项
                        # 注意：这里的值仅作为标记使用，不使用其数值
                        if pd.notna(value) and value != '' and value != 0:
                            scores.append(ScoreItem(
                                question_name=str(knowledge_point),
                                deduction=1.0,  # 固定值，仅表示有扣分，不代表真实扣分
                                category=FileService._get_question_category(str(knowledge_point))
                            ))
                            logger.info(f"{student_name} 在知识点 '{knowledge_point}' 有扣分标记")
                    
                    # 添加学生成绩（即使没有扣分也添加，表示满分）
                    student_scores.append(StudentScore(
                        student_name=student_name,
                        scores=scores,
                        total_score=total_score
                    ))
                    logger.info(f"成功添加 {student_name} 的成绩记录，共 {len(scores)} 个扣分知识点")
                
                except Exception as e:
                    logger.error(f"处理学生行时出错: {str(e)}")
                    continue
            
            logger.info(f"Excel处理完成，共处理 {len(student_scores)} 个学生的成绩")
            return student_scores
        except Exception as e:
            error_msg = f"处理Excel文件时出错: {str(e)}"
            logger.error(error_msg)
            raise Exception(error_msg)

    @staticmethod
    async def process_word(file_path: str) -> List[StudentScore]:
        """处理Word文件并提取成绩信息"""
        try:
            doc = Document(file_path)
            student_scores = []
            current_student = None
            current_scores = []
            
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                    
                # 尝试识别学生姓名（假设学生姓名是单独的一行）
                if not current_student and len(text.split()) <= 2:
                    current_student = text
                    continue
                
                # 尝试识别成绩信息
                if current_student:
                    try:
                        # 假设成绩格式为：题目名称：扣分
                        if '：' in text or ':' in text:
                            parts = text.replace('：', ':').split(':')
                            if len(parts) == 2:
                                question_name = parts[0].strip()
                                deduction = float(parts[1].strip())
                                category = FileService._get_question_category(question_name)
                                current_scores.append(ScoreItem(
                                    question_name=question_name,
                                    deduction=deduction,
                                    category=category
                                ))
                    except (ValueError, TypeError):
                        continue
            
            if current_student and current_scores:
                total_score = sum(score.deduction for score in current_scores)
                student_scores.append(StudentScore(
                    student_name=current_student,
                    scores=current_scores,
                    total_score=total_score
                ))
            
            return student_scores
        except Exception as e:
            raise Exception(f"处理Word文件时出错: {str(e)}")

    @staticmethod
    async def process_ppt(file_path: str) -> List[StudentScore]:
        """处理PPT文件并提取成绩信息"""
        try:
            prs = Presentation(file_path)
            student_scores = []
            
            for slide in prs.slides:
                student_name = None
                scores = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if not text:
                            continue
                            
                        # 尝试识别学生姓名
                        if not student_name and len(text.split()) <= 2:
                            student_name = text
                            continue
                        
                        # 尝试识别成绩信息
                        if student_name:
                            try:
                                if '：' in text or ':' in text:
                                    parts = text.replace('：', ':').split(':')
                                    if len(parts) == 2:
                                        question_name = parts[0].strip()
                                        deduction = float(parts[1].strip())
                                        category = FileService._get_question_category(question_name)
                                        scores.append(ScoreItem(
                                            question_name=question_name,
                                            deduction=deduction,
                                            category=category
                                        ))
                            except (ValueError, TypeError):
                                continue
                
                if student_name and scores:
                    total_score = sum(score.deduction for score in scores)
                    student_scores.append(StudentScore(
                        student_name=student_name,
                        scores=scores,
                        total_score=total_score
                    ))
            
            return student_scores
        except Exception as e:
            raise Exception(f"处理PPT文件时出错: {str(e)}")

    @staticmethod
    def _get_question_category(question_name: str) -> str:
        """根据题目名称判断题目类型"""
        question_name_lower = question_name.lower()
        
        # 根据关键词判断类别
        if '选择' in question_name or '单选' in question_name or '多选' in question_name:
            return '选择题'
        elif '填空' in question_name:
            return '填空题'
        elif '计算' in question_name:
            return '计算题'
        elif '应用' in question_name or '解答' in question_name:
            return '应用题'
        elif '巧' in question_name or '巧题' in question_name:
            return '巧题'
        elif '判断' in question_name:
            return '判断题'
        elif '简答' in question_name:
            return '简答题'
        elif '作图' in question_name or '画图' in question_name:
            return '作图题'
        else:
            # 如果无法识别，返回原名称作为类别
            return question_name 